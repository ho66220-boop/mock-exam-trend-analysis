from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TABLE_DIR = ROOT / "output" / "tables"
REPORT_DIR = ROOT / "output" / "reports"
PPT_PATH = REPORT_DIR / "final_insight_report_internal_share.pptx"

FONT = "Malgun Gothic"
TITLE = RGBColor(31, 45, 61)
TEXT = RGBColor(55, 65, 81)
MUTED = RGBColor(107, 114, 128)
BLUE = RGBColor(47, 128, 237)
GREEN = RGBColor(39, 174, 96)
ORANGE = RGBColor(242, 153, 74)
RED = RGBColor(235, 87, 87)
PURPLE = RGBColor(155, 81, 224)
GRAY = RGBColor(130, 130, 130)


def read_csv(name: str) -> pd.DataFrame:
    return pd.read_csv(TABLE_DIR / name)


def fmt(value, digits=1):
    if pd.isna(value):
        return "-"
    if isinstance(value, float):
        return f"{value:.{digits}f}"
    return str(value)


def fmt_pct(value):
    if pd.isna(value):
        return "-"
    return f"{value:.1%}"


def add_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    return slide


def set_text(shape, text, size=16, bold=False, color=TEXT):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color


def title(slide, text, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.3), Inches(12.2), Inches(0.55))
    set_text(box, text, 24, True, TITLE)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.88), Inches(12.2), Inches(0.35))
        set_text(sub, subtitle, 10, False, MUTED)


def bullets(slide, items, x=0.75, y=1.35, w=12.0, h=5.4, size=16):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(8)


def note(slide, text):
    box = slide.shapes.add_textbox(Inches(0.65), Inches(6.85), Inches(12.0), Inches(0.28))
    set_text(box, text, 8, False, MUTED)


def table(slide, df, columns, headers, x, y, w, h, size=8):
    shape = slide.shapes.add_table(len(df) + 1, len(columns), Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = shape.table
    for i, header in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(230, 236, 245)
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT
            p.font.size = Pt(size)
            p.font.bold = True
            p.font.color.rgb = TITLE
            p.alignment = PP_ALIGN.CENTER
    for r, (_, row) in enumerate(df.iterrows(), 1):
        for c, col in enumerate(columns):
            value = row[col]
            text = fmt_pct(value) if "rate" in col else fmt(value, 3 if col in {"pearson_corr", "cv_r2_mean"} else 1)
            cell = tbl.cell(r, c)
            cell.text = text
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT
                p.font.size = Pt(size)
                p.font.color.rgb = TEXT
                p.alignment = PP_ALIGN.CENTER


def style_chart(chart, chart_title, x_title, y_title, labels=False):
    chart.has_title = True
    chart.chart_title.text_frame.text = chart_title
    chart.chart_title.text_frame.paragraphs[0].font.name = FONT
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.category_axis.has_title = True
    chart.category_axis.axis_title.text_frame.text = x_title
    chart.category_axis.axis_title.text_frame.paragraphs[0].font.name = FONT
    chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(9)
    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.text = y_title
    chart.value_axis.axis_title.text_frame.paragraphs[0].font.name = FONT
    chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(9)
    chart.value_axis.has_major_gridlines = True
    if labels:
        chart.plots[0].has_data_labels = True
        chart.plots[0].data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        chart.plots[0].data_labels.font.size = Pt(8)
        chart.plots[0].data_labels.font.name = FONT


def line_chart(slide, categories, series, x, y, w, h, chart_title, x_title, y_title):
    data = CategoryChartData()
    data.categories = categories
    for name, values in series.items():
        data.add_series(name, values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(x), Inches(y), Inches(w), Inches(h), data).chart
    style_chart(chart, chart_title, x_title, y_title)
    for idx, color in enumerate([BLUE, GREEN, ORANGE, RED]):
        if idx < len(chart.series):
            chart.series[idx].format.line.color.rgb = color


def bar_chart(slide, categories, series, x, y, w, h, chart_title, x_title, y_title):
    data = CategoryChartData()
    data.categories = categories
    for name, values in series.items():
        data.add_series(name, values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data).chart
    style_chart(chart, chart_title, x_title, y_title, True)
    for idx, color in enumerate([BLUE, GREEN, ORANGE, RED, PURPLE, GRAY]):
        if idx < len(chart.series):
            chart.series[idx].format.fill.solid()
            chart.series[idx].format.fill.fore_color.rgb = color


def hbar_chart(slide, categories, values, x, y, w, h, chart_title, x_title, y_title):
    data = CategoryChartData()
    data.categories = categories
    data.add_series("값", values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data).chart
    style_chart(chart, chart_title, y_title, x_title, True)
    chart.has_legend = False
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = BLUE


def build():
    monthly = read_csv("monthly_level_percentile_trend.csv")
    level = read_csv("segment_level_summary.csv")
    strength = read_csv("segment_level_strength_summary.csv")
    participation = read_csv("segment_level_participation_summary.csv")
    volatility = read_csv("segment_level_volatility_summary.csv")
    recommendation = read_csv("segment_recommendation_table.csv")
    mobility_strength = read_csv("mobility_by_strength.csv")
    mobility_participation = read_csv("mobility_by_participation.csv")
    mobility_volatility = read_csv("mobility_by_volatility.csv")
    mobility_detail = read_csv("mobility_by_detailed_segment.csv")
    mobility_profiles = read_csv("student_mobility_profiles.csv")
    inquiry_compare = read_csv("inquiry_switch_vs_keep_summary.csv")
    inquiry_timing = read_csv("inquiry_switch_timing_summary.csv")
    corr = read_csv("subject_feature_correlations.csv").sort_values("abs_corr", ascending=False).head(8)
    gap = read_csv("exam_to_csat_gap_summary.csv")
    model = read_csv("model_comparison.csv")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    s = add_slide(prs)
    box = s.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.8), Inches(0.9))
    set_text(box, "샘플 재원생 모의고사-수능 결과 인사이트 종합 리포트", 30, True, TITLE)
    set_text(s.shapes.add_textbox(Inches(0.85), Inches(2.25), Inches(11.7), Inches(0.5)), "공개 데모용 | 더미 데이터 기반 예시 자료", 16, False, TEXT)
    bullets(s, ["개인별 수능 성적을 단정적으로 예측하지 않습니다.", "샘플 재원생의 학습 패턴과 수능 결과의 관계를 상담 참고 신호로 정리합니다.", "모든 차트와 표는 PPT에서 직접 수정 가능한 객체입니다."], y=3.25, size=17)

    s = add_slide(prs)
    title(s, "데이터 기준")
    bullets(s, ["분석 대상: 수능 이전 모의고사 기록과 수능 결과가 모두 있는 116명", "핵심 성과 지표: 국어, 수학, 탐구1, 탐구2 백분위 평균", "영어는 등급 자료이므로 핵심 평균에는 포함하지 않고 보조 지표로 활용", "점수 0은 미응시 또는 미기록 가능성이 높아 결측값 처리", "상위권/중위권/하위권은 수능 이전 모의고사 핵심 평균의 3분위 기준"], size=17)

    s = add_slide(prs)
    title(s, "월별 상/중/하위권 평균 백분위 흐름")
    pivot = monthly.pivot(index="month", columns="pre_level", values="avg_core_percentile").reset_index()
    pivot["order"] = pivot["month"].str.replace("월", "", regex=False).astype(int)
    pivot = pivot.sort_values("order")
    line_chart(s, pivot["month"].tolist(), {"상위권": pivot["상위권"].round(1).tolist(), "중위권": pivot["중위권"].round(1).tolist(), "하위권": pivot["하위권"].round(1).tolist()}, 0.65, 1.15, 8.0, 4.8, "월별 평균 백분위", "시험 월", "국어/수학/탐구 평균 백분위")
    table(s, pivot[["month", "상위권", "중위권", "하위권"]], ["month", "상위권", "중위권", "하위권"], ["월", "상위권", "중위권", "하위권"], 8.9, 1.35, 3.7, 4.1, 8)

    order = {"상위권": 0, "중위권": 1, "하위권": 2}
    level = level.sort_values("pre_level", key=lambda x: x.map(order))
    s = add_slide(prs)
    title(s, "전체 요약")
    bar_chart(s, level["pre_level"].tolist(), {"수능 이전 평균": level["avg_pre_core_mean"].round(1).tolist(), "수능 평균": level["avg_csat_core_mean"].round(1).tolist()}, 0.65, 1.15, 7.5, 4.8, "권역별 수능 이전 평균과 수능 평균", "권역", "핵심 평균 백분위")
    table(s, level, ["pre_level", "n", "top_25pct_rate", "rise_rate", "drop_rate"], ["구분", "학생 수", "상위 25%", "상승률", "하락률"], 8.35, 1.55, 4.2, 2.2, 9)
    note(s, "하위권은 상승률이 존재하더라도 상위권 진입 비율은 낮아 기본 성취 수준 관리가 중요합니다.")

    s = add_slide(prs)
    title(s, "과목 강점 유형별 인사이트")
    levels = ["상위권", "중위권", "하위권"]
    strengths = ["균형형", "국어형", "수학형", "탐구형"]
    sp = strength.pivot(index="pre_level", columns="strength_label", values="avg_csat_core_mean").reindex(levels).reindex(columns=strengths)
    bar_chart(s, levels, {c: sp[c].fillna(0).round(1).tolist() for c in strengths}, 0.65, 1.15, 12.0, 4.9, "성취권 x 과목 강점 유형별 수능 평균", "권역", "수능 핵심 평균")
    note(s, "상위권에서는 균형형 학생의 수능 결과가 가장 안정적으로 높았습니다.")

    s = add_slide(prs)
    title(s, "과목 강점 유형별 상세표")
    table(s, strength.sort_values(["pre_level", "avg_csat_core_mean"], ascending=[True, False]), ["pre_level", "strength_label", "n", "avg_csat_core_mean", "rise_rate", "drop_rate"], ["권역", "강점", "학생 수", "수능 평균", "상승률", "하락률"], 0.55, 1.05, 12.2, 5.8, 7)

    s = add_slide(prs)
    title(s, "응시 횟수별 인사이트")
    part_groups = ["적게 응시(1-3회)", "보통 응시(4-6회)", "많이 응시(7회 이상)"]
    pp = participation.pivot(index="pre_level", columns="participation_group", values="avg_csat_core_mean").reindex(levels).reindex(columns=part_groups)
    bar_chart(s, levels, {c: pp[c].fillna(0).round(1).tolist() for c in part_groups}, 0.65, 1.15, 8.1, 4.8, "성취권 x 응시 횟수별 수능 평균", "권역", "수능 핵심 평균")
    table(s, participation.sort_values(["pre_level", "avg_csat_core_mean"], ascending=[True, False]), ["pre_level", "participation_group", "n", "avg_csat_core_mean", "top_25pct_rate"], ["권역", "응시", "학생 수", "수능 평균", "상위 25%"], 8.95, 1.25, 3.85, 4.6, 7)

    s = add_slide(prs)
    title(s, "변동성별 인사이트")
    vp = volatility.pivot(index="pre_level", columns="volatility_group", values="avg_csat_core_mean").reindex(levels).reindex(columns=["변동성 작음", "변동성 큼"])
    bar_chart(s, levels, {"변동성 작음": vp["변동성 작음"].round(1).tolist(), "변동성 큼": vp["변동성 큼"].round(1).tolist()}, 0.65, 1.15, 7.2, 4.8, "성취권 x 변동성별 수능 평균", "권역", "수능 핵심 평균")
    table(s, volatility.sort_values(["pre_level", "volatility_group"]), ["pre_level", "volatility_group", "n", "avg_csat_core_mean", "avg_csat_minus_pre", "drop_rate"], ["권역", "변동성", "학생 수", "수능 평균", "수능-이전", "하락률"], 8.05, 1.25, 4.65, 4.0, 7)

    s = add_slide(prs)
    title(s, "계층 이동 분석")
    counts = mobility_profiles["mobility_type"].value_counts()
    bullets(s, [f"상향 이동: {int(counts.get('상향 이동', 0))}명 ({counts.get('상향 이동', 0)/len(mobility_profiles):.1%})", f"유지: {int(counts.get('유지', 0))}명 ({counts.get('유지', 0)/len(mobility_profiles):.1%})", f"하향 이동: {int(counts.get('하향 이동', 0))}명 ({counts.get('하향 이동', 0)/len(mobility_profiles):.1%})"], x=0.8, y=1.25, w=4.0, h=1.8, size=18)
    ms = mobility_strength[mobility_strength["n"] >= 5].sort_values("up_rate")
    hbar_chart(s, ms["strength_label"].tolist(), (ms["up_rate"] * 100).round(1).tolist(), 5.05, 1.15, 7.4, 3.6, "강점 유형별 상향 이동률", "상향 이동률(%)", "강점 유형")
    table(s, mobility_participation, ["participation_group", "n", "up_rate", "down_rate", "avg_level_move"], ["응시", "학생 수", "상향", "하향", "평균 이동"], 0.8, 4.0, 5.2, 1.7, 8)
    table(s, mobility_volatility, ["volatility_group", "n", "up_rate", "down_rate", "avg_level_move"], ["변동성", "학생 수", "상향", "하향", "평균 이동"], 6.35, 4.0, 5.5, 1.2, 8)

    s = add_slide(prs)
    title(s, "계층 이동이 활발한 세부 유형")
    table(s, mobility_detail[mobility_detail["n"] >= 3].head(8), ["pre_level", "strength_label", "participation_group", "volatility_group", "n", "up_rate", "down_rate", "avg_level_move"], ["시작", "강점", "응시", "변동성", "학생 수", "상향", "하향", "평균 이동"], 0.45, 1.05, 12.5, 4.5, 7)
    note(s, "표본 수가 작은 조합은 '작년 데이터에서 관찰된 사례'로만 표현하는 것이 안전합니다.")

    s = add_slide(prs)
    title(s, "전달용 핵심 세그먼트")
    table(s, recommendation.head(10), ["pre_level", "strength_label", "participation_group", "volatility_group", "n", "avg_csat_core_mean", "insight_message"], ["권역", "강점", "응시", "변동성", "학생 수", "수능 평균", "전달 메시지"], 0.35, 0.95, 12.65, 5.9, 6)

    s = add_slide(prs)
    title(s, "탐구 과목 변경 인사이트")
    bar_chart(s, inquiry_compare["group"].tolist(), {"수능 이전 탐구 평균": inquiry_compare["avg_pre_all_inquiry_mean"].round(1).tolist(), "수능 탐구 평균": inquiry_compare["avg_csat_inquiry_mean"].round(1).tolist()}, 0.65, 1.15, 5.7, 4.6, "탐구 변경 여부별 탐구 평균", "구분", "탐구 평균 백분위")
    table(s, inquiry_timing, ["change_month", "n", "avg_post_minus_pre", "avg_csat_minus_pre_change", "avg_csat_minus_post_change", "benefit_rate_csat"], ["변경", "n", "후-전", "수능-전", "수능-후", "수능 이득"], 6.65, 1.2, 6.0, 3.5, 7)
    note(s, "탐구 변경 학생은 17명으로 표본이 작아 시기별 결론은 참고 신호로만 해석합니다.")

    s = add_slide(prs)
    title(s, "모의고사 지표와 수능 결과의 관계")
    table(s, corr, ["label", "feature", "n", "pearson_corr"], ["과목", "수능 이전 지표", "표본 수", "상관계수"], 0.65, 1.15, 6.0, 4.2, 8)
    bullets(s, ["수능 결과와 가장 강하게 연결된 지표는 특정 시험 한 번보다 수능 이전 평균과 최근 평균이었습니다.", "특히 수학 평균, 국어 평균 지표가 강하게 나타났습니다."], x=7.1, y=1.5, w=5.5, h=2.0, size=16)

    s = add_slide(prs)
    title(s, "시험별 수능과의 차이")
    best_rows = []
    for _, group in gap.groupby("label"):
        cand = group[group["n"] >= 30].sort_values("mean_abs_gap").head(1)
        if not cand.empty:
            best_rows.append(cand.iloc[0])
    table(s, pd.DataFrame(best_rows), ["label", "exam_name", "n", "mean_abs_gap", "pearson_corr"], ["과목", "수능과 가까운 시험", "표본 수", "평균 절대 차이", "상관계수"], 0.55, 1.2, 12.1, 3.2, 8)

    s = add_slide(prs)
    title(s, "모델 설명력 비교")
    core = model[model["target"] == "csat_core_mean"].sort_values("cv_r2_mean")
    hbar_chart(s, core["model"].tolist(), core["cv_r2_mean"].round(3).tolist(), 0.65, 1.2, 6.8, 4.5, "수능 핵심 평균 모델 설명력", "5-fold R2", "모델")
    table(s, core.sort_values("cv_r2_mean", ascending=False).head(6), ["model", "n", "cv_r2_mean", "cv_mae_mean", "cv_rmse_mean"], ["모델", "표본", "R2", "MAE", "RMSE"], 7.85, 1.35, 4.6, 3.2, 8)
    note(s, "개인별 예측기가 아니라 현재 데이터에서 어떤 설명 방식이 타당한지 확인하는 비교입니다.")

    s = add_slide(prs)
    title(s, "학생에게 전달할 때의 표현 예시")
    bullets(s, ["샘플 데이터에서는 상위권 학생 중 국어, 수학, 탐구가 고르게 받쳐주는 균형형 학생의 수능 결과가 가장 안정적으로 높았습니다.", "월별 평균 백분위 흐름을 보면 권역별 격차와 유지 양상이 드러나므로, 학생에게 현재 위치와 변화 방향을 설명하는 데 활용할 수 있습니다.", "샘플 데이터에서는 같은 계층을 유지한 학생이 가장 많고, 일부 상향/하향 이동 사례가 관찰됩니다.", "상위권이라도 월별 성적 변동이 큰 학생은 수능에서 하락한 사례가 더 많이 관찰되어, 흔들림 관리가 중요합니다.", "탐구 과목 변경은 변경 자체의 효과보다 변경 전후에 충분한 검증 기록을 확보했는지가 더 중요해 보입니다."], y=1.25, size=15)

    s = add_slide(prs)
    title(s, "해석 주의")
    bullets(s, ["본 분석은 샘플 데이터의 패턴 요약이며 개인별 결과를 보장하지 않습니다.", "표본 수가 작은 세그먼트는 방향성 참고용으로만 사용해야 합니다.", "응시 횟수와 수능 결과의 관계는 인과가 아니라 상관 또는 집단 특성의 반영일 수 있습니다.", "계층 이동은 3분위 기준의 상대적 이동이므로 경계 부근 학생은 작은 점수 변화로도 이동할 수 있습니다.", "탐구 변경 분석은 표본 수가 작아 시기별 결론을 강하게 일반화하면 안 됩니다."], y=1.35, size=17)

    REPORT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(PPT_PATH)
    print(PPT_PATH)


if __name__ == "__main__":
    build()
