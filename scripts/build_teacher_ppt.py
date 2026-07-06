from pathlib import Path

import pandas as pd

import exam_meta
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TABLE_DIR = ROOT / "output" / "tables"
REPORT_DIR = ROOT / "output" / "reports"
PPT_PATH = REPORT_DIR / "mock_exam_insight_teacher_briefing.pptx"

FONT = "Malgun Gothic"
TITLE_COLOR = RGBColor(31, 45, 61)
TEXT_COLOR = RGBColor(55, 65, 81)
BLUE = RGBColor(47, 128, 237)
GREEN = RGBColor(39, 174, 96)
ORANGE = RGBColor(242, 153, 74)
RED = RGBColor(235, 87, 87)
PURPLE = RGBColor(155, 81, 224)
GRAY = RGBColor(130, 130, 130)
LIGHT_BG = RGBColor(247, 249, 252)


def read_csv(name: str) -> pd.DataFrame:
    path = TABLE_DIR / name
    if not path.exists():
        raise FileNotFoundError(f"Missing {path}")
    return pd.read_csv(path)


def fmt(value, digits: int = 1) -> str:
    if pd.isna(value):
        return "-"
    if isinstance(value, float):
        return f"{value:.{digits}f}"
    return str(value)


def fmt_pct(value) -> str:
    if pd.isna(value):
        return "-"
    return f"{value:.1%}"


def set_text(shape, text: str, size: int = 18, bold: bool = False, color=TEXT_COLOR) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.55))
    set_text(title_box, title, size=24, bold=True, color=TITLE_COLOR)
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.58), Inches(0.92), Inches(12.1), Inches(0.35))
        set_text(subtitle_box, subtitle, size=11, color=GRAY)


def add_note(slide, text: str, y: float = 6.85) -> None:
    box = slide.shapes.add_textbox(Inches(0.65), Inches(y), Inches(12.0), Inches(0.3))
    set_text(box, text, size=9, color=GRAY)


def add_bullets(slide, items: list[str], x=0.75, y=1.55, w=12.0, h=4.8, size=17) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(8)


def add_table(slide, df: pd.DataFrame, headers: list[str], columns: list[str], x, y, w, h, font_size=9) -> None:
    rows = len(df) + 1
    cols = len(columns)
    table_shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table

    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(230, 236, 245)
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = TITLE_COLOR
            p.alignment = PP_ALIGN.CENTER

    for row_idx, (_, row) in enumerate(df.iterrows(), start=1):
        for col_idx, column in enumerate(columns):
            value = row[column]
            if "rate" in column:
                text = fmt_pct(value)
            else:
                text = fmt(value)
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT
                p.font.size = Pt(font_size)
                p.font.color.rgb = TEXT_COLOR
                p.alignment = PP_ALIGN.CENTER


def style_chart(chart, title: str, category_axis: str, value_axis: str, show_data_labels=False) -> None:
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].font.name = FONT
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    chart.category_axis.has_title = True
    chart.category_axis.axis_title.text_frame.text = category_axis
    chart.category_axis.axis_title.text_frame.paragraphs[0].font.name = FONT
    chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)

    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.text = value_axis
    chart.value_axis.axis_title.text_frame.paragraphs[0].font.name = FONT
    chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
    chart.value_axis.has_major_gridlines = True

    if show_data_labels:
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        plot.data_labels.font.size = Pt(8)
        plot.data_labels.font.name = FONT


def add_line_chart(slide, categories, series: dict[str, list[float]], x, y, w, h, title, x_title, y_title) -> None:
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series.items():
        chart_data.add_series(name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, Inches(x), Inches(y), Inches(w), Inches(h), chart_data
    ).chart
    style_chart(chart, title, x_title, y_title)
    for idx, color in enumerate([BLUE, GREEN, ORANGE, RED, PURPLE]):
        if idx < len(chart.series):
            chart.series[idx].format.line.color.rgb = color
            chart.series[idx].marker.format.fill.solid()
            chart.series[idx].marker.format.fill.fore_color.rgb = color


def add_bar_chart(slide, categories, series: dict[str, list[float]], x, y, w, h, title, x_title, y_title) -> None:
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series.items():
        chart_data.add_series(name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), chart_data
    ).chart
    style_chart(chart, title, x_title, y_title, show_data_labels=True)
    colors = [BLUE, GREEN, ORANGE, RED, PURPLE, GRAY]
    for idx, serie in enumerate(chart.series):
        serie.format.fill.solid()
        serie.format.fill.fore_color.rgb = colors[idx % len(colors)]


def add_horizontal_bar(slide, categories, values, x, y, w, h, title, x_title, y_title, series_name="값") -> None:
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), chart_data
    ).chart
    style_chart(chart, title, y_title, x_title, show_data_labels=True)
    chart.has_legend = False
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = BLUE


def add_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    return slide


def build_presentation() -> None:
    monthly = read_csv("monthly_level_percentile_trend.csv")
    level = read_csv("segment_level_summary.csv")
    strength = read_csv("segment_level_strength_summary.csv")
    participation = read_csv("segment_level_participation_summary.csv")
    volatility = read_csv("segment_level_volatility_summary.csv")
    mobility_strength = read_csv("mobility_by_strength.csv")
    mobility_participation = read_csv("mobility_by_participation.csv")
    inquiry_compare = read_csv("inquiry_switch_vs_keep_summary.csv")
    inquiry_timing = read_csv("inquiry_switch_timing_summary.csv")
    model = read_csv("model_comparison.csv")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Title
    slide = add_slide(prs)
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.8), Inches(0.9))
    set_text(title_box, "샘플 재원생 모의고사-수능 결과 인사이트", size=32, bold=True, color=TITLE_COLOR)
    sub_box = slide.shapes.add_textbox(Inches(0.85), Inches(2.35), Inches(11.4), Inches(0.8))
    set_text(
        sub_box,
        "공개 데모용 | 개인별 예측이 아닌 상담 참고 패턴 요약",
        size=17,
        color=TEXT_COLOR,
    )
    add_bullets(
        slide,
        [
            "분석 대상: 수능 이전 모의고사 기록과 수능 결과가 모두 있는 116명",
            "핵심 지표: 국어/수학/탐구1/탐구2 백분위 평균",
            "원본 데이터와 중간 CSV는 비공개, 리포트와 시각화만 공유",
        ],
        x=0.95,
        y=3.45,
        w=11.2,
        h=2.2,
        size=16,
    )

    # 2. Key messages
    slide = add_slide(prs)
    add_title(slide, "핵심 메시지", "학생에게는 '예측'보다 '샘플 데이터 패턴'으로 전달")
    add_bullets(
        slide,
        [
            "상위권은 균형형 학생의 수능 결과가 가장 안정적으로 높았습니다.",
            "중위권과 하위권은 일부 상향 이동 사례가 있으며, 약점 축 관리가 중요합니다.",
            "응시 횟수는 인과가 아니지만, 4-6회 보통 응시 그룹이 비교적 안정적이었습니다.",
            "상위권이라도 변동성이 큰 학생은 하락 위험을 함께 점검해야 합니다.",
            "탐구 변경은 '바꾸면 오른다'가 아니라, 변경 후 충분히 검증할 시간이 있는지가 핵심입니다.",
        ],
        y=1.55,
        size=18,
    )

    # 3. Monthly trend
    slide = add_slide(prs)
    add_title(slide, "월별 상/중/하위권 평균 백분위 흐름", "차트의 축 제목과 계열명은 PowerPoint에서 직접 수정 가능")
    pivot = monthly.pivot(index="month", columns="pre_level", values="avg_core_percentile").reset_index()
    pivot["month_order"] = pivot["month"].map(exam_meta.label_sort_key)
    pivot = pivot.sort_values("month_order")
    add_line_chart(
        slide,
        pivot["month"].tolist(),
        {
            "상위권": pivot["상위권"].round(1).tolist(),
            "중위권": pivot["중위권"].round(1).tolist(),
            "하위권": pivot["하위권"].round(1).tolist(),
        },
        0.7,
        1.35,
        12.0,
        5.2,
        "월별 평균 백분위",
        "시험 월",
        "국어/수학/탐구 평균 백분위",
    )
    add_note(slide, "상/중/하위권은 수능 이전 모의고사 핵심 평균의 3분위 기준입니다.")

    # 4. Level summary
    slide = add_slide(prs)
    add_title(slide, "상/중/하위권별 수능 결과 요약")
    order = {"상위권": 0, "중위권": 1, "하위권": 2}
    level = level.sort_values("pre_level", key=lambda s: s.map(order))
    add_bar_chart(
        slide,
        level["pre_level"].tolist(),
        {
            "수능 이전 평균": level["avg_pre_core_mean"].round(1).tolist(),
            "수능 평균": level["avg_csat_core_mean"].round(1).tolist(),
        },
        0.65,
        1.25,
        7.35,
        4.8,
        "수능 이전 평균과 수능 평균 비교",
        "권역",
        "핵심 평균 백분위",
    )
    table_df = level[["pre_level", "n", "top_25pct_rate", "rise_rate", "drop_rate"]].copy()
    add_table(
        slide,
        table_df,
        ["권역", "학생 수", "상위 25%", "상승률", "하락률"],
        ["pre_level", "n", "top_25pct_rate", "rise_rate", "drop_rate"],
        8.25,
        1.65,
        4.3,
        2.1,
        font_size=10,
    )
    add_note(slide, "상위권은 높은 결과를 유지했고, 중위권/하위권은 일부 상승 사례가 확인됩니다.")

    # 5. Strength by level
    slide = add_slide(prs)
    add_title(slide, "성취권 x 과목 강점 유형별 수능 결과")
    strengths = ["균형형", "국어형", "수학형", "탐구형"]
    levels = ["상위권", "중위권", "하위권"]
    strength_pivot = (
        strength.pivot(index="pre_level", columns="strength_label", values="avg_csat_core_mean")
        .reindex(levels)
        .reindex(columns=strengths)
    )
    add_bar_chart(
        slide,
        levels,
        {name: strength_pivot[name].round(1).fillna(0).tolist() for name in strengths},
        0.65,
        1.25,
        12.0,
        5.2,
        "권역별 강점 유형 수능 평균",
        "권역",
        "수능 핵심 평균",
    )
    add_note(slide, "상위권은 균형형이 가장 안정적이고, 하위권은 수학형의 상승 가능성이 상대적으로 두드러졌습니다.")

    # 6. Participation
    slide = add_slide(prs)
    add_title(slide, "응시 횟수별 수능 결과")
    part_groups = ["적게 응시(1-3회)", "보통 응시(4-6회)", "많이 응시(7회 이상)"]
    part_pivot = (
        participation.pivot(index="pre_level", columns="participation_group", values="avg_csat_core_mean")
        .reindex(levels)
        .reindex(columns=part_groups)
    )
    add_bar_chart(
        slide,
        levels,
        {name: part_pivot[name].round(1).fillna(0).tolist() for name in part_groups},
        0.65,
        1.25,
        12.0,
        5.1,
        "권역별 응시 횟수 수능 평균",
        "권역",
        "수능 핵심 평균",
    )
    add_note(slide, "응시 횟수는 인과가 아니라 재원 기간, 성실성, 시험 참여 성향이 섞인 지표로 해석합니다.")

    # 7. Volatility
    slide = add_slide(prs)
    add_title(slide, "변동성별 수능 결과와 하락 위험")
    vol_pivot = (
        volatility.pivot(index="pre_level", columns="volatility_group", values="avg_csat_core_mean")
        .reindex(levels)
        .reindex(columns=["변동성 작음", "변동성 큼"])
    )
    add_bar_chart(
        slide,
        levels,
        {
            "변동성 작음": vol_pivot["변동성 작음"].round(1).tolist(),
            "변동성 큼": vol_pivot["변동성 큼"].round(1).tolist(),
        },
        0.65,
        1.25,
        7.1,
        4.9,
        "권역별 변동성 수능 평균",
        "권역",
        "수능 핵심 평균",
    )
    add_table(
        slide,
        volatility.sort_values(["pre_level", "volatility_group"]),
        ["권역", "변동성", "학생 수", "하락률"],
        ["pre_level", "volatility_group", "n", "drop_rate"],
        8.05,
        1.45,
        4.7,
        3.2,
        font_size=9,
    )
    add_note(slide, "상위권에서는 변동성이 큰 학생의 하락률이 더 높게 관찰됩니다.")

    # 8. Mobility
    slide = add_slide(prs)
    add_title(slide, "계층 이동 분석")
    ms = mobility_strength[mobility_strength["n"] >= 5].sort_values("up_rate")
    add_horizontal_bar(
        slide,
        ms["strength_label"].tolist(),
        (ms["up_rate"] * 100).round(1).tolist(),
        0.75,
        1.25,
        5.9,
        4.7,
        "강점 유형별 상향 이동률",
        "상향 이동률(%)",
        "강점 유형",
        "상향 이동률",
    )
    add_table(
        slide,
        mobility_participation,
        ["응시 횟수", "학생 수", "상향", "하향"],
        ["participation_group", "n", "up_rate", "down_rate"],
        7.0,
        1.65,
        5.6,
        1.7,
        font_size=10,
    )
    add_bullets(
        slide,
        [
            "전체 상향 이동: 16명(13.8%)",
            "전체 유지: 85명(73.3%)",
            "전체 하향 이동: 15명(12.9%)",
        ],
        x=7.05,
        y=4.0,
        w=5.5,
        h=1.5,
        size=15,
    )

    # 9. Inquiry switching
    slide = add_slide(prs)
    add_title(slide, "탐구 과목 변경 분석")
    add_bar_chart(
        slide,
        inquiry_compare["group"].tolist(),
        {
            "수능 이전 탐구 평균": inquiry_compare["avg_pre_all_inquiry_mean"].round(1).tolist(),
            "수능 탐구 평균": inquiry_compare["avg_csat_inquiry_mean"].round(1).tolist(),
        },
        0.65,
        1.25,
        5.9,
        4.8,
        "탐구 변경 여부별 탐구 평균",
        "구분",
        "탐구 평균 백분위",
    )
    add_table(
        slide,
        inquiry_timing,
        ["시기", "n", "변경 후-전", "수능-변경 전", "수능 이득률"],
        ["change_month", "n", "avg_post_minus_pre", "avg_csat_minus_pre_change", "benefit_rate_csat"],
        6.85,
        1.35,
        5.95,
        3.35,
        font_size=8,
    )
    add_note(slide, "탐구 변경 학생은 17명으로 표본이 작아, 시기별 결과는 참고 신호로만 해석합니다.")

    # 10. Model comparison
    slide = add_slide(prs)
    add_title(slide, "모델 설명력 비교", "개인별 예측기가 아니라 패턴 설명 방식 검증")
    core = model[model["target"] == "csat_core_mean"].sort_values("cv_r2_mean")
    add_horizontal_bar(
        slide,
        core["model"].tolist(),
        core["cv_r2_mean"].round(3).tolist(),
        0.75,
        1.25,
        6.5,
        4.9,
        "수능 핵심 평균 설명력",
        "5-fold 교차검증 R2",
        "모델",
        "R2",
    )
    add_table(
        slide,
        core.sort_values("cv_r2_mean", ascending=False)[["model", "cv_r2_mean", "cv_mae_mean"]].head(4),
        ["모델", "R2", "MAE"],
        ["model", "cv_r2_mean", "cv_mae_mean"],
        7.7,
        1.55,
        4.7,
        2.0,
        font_size=10,
    )
    add_bullets(
        slide,
        [
            "복잡한 모델이 압도적으로 앞서지 않았습니다.",
            "상담용 설명에는 선형 계열 지표가 더 안전합니다.",
            "수능 이전 평균, 수학 평균, 국어 평균이 핵심 설명 변수입니다.",
        ],
        x=7.7,
        y=4.05,
        w=4.9,
        h=1.8,
        size=14,
    )

    # 11. Closing
    slide = add_slide(prs)
    add_title(slide, "담임 상담 활용 포인트")
    add_bullets(
        slide,
        [
            "학생에게는 '성적 예측'이 아니라 '샘플 데이터 패턴'으로 설명합니다.",
            "현재 위치는 월별 상/중/하위권 평균 흐름과 비교해 설명합니다.",
            "상위권은 균형 유지와 변동성 관리가 핵심입니다.",
            "중위권은 약점 축 보완과 계층 상향 가능성을 함께 봅니다.",
            "탐구 변경은 변경 후 검증 기록이 충분한지 반드시 확인합니다.",
        ],
        y=1.45,
        size=18,
    )
    add_note(slide, "모든 차트와 표는 PowerPoint에서 직접 수정 가능한 객체로 생성했습니다.", y=6.55)

    PPT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(PPT_PATH)
    print(PPT_PATH)


if __name__ == "__main__":
    build_presentation()
